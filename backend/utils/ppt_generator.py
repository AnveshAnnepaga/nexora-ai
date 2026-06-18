import os
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from typing import Dict, Any

def create_startup_pitch_deck(state: Dict[str, Any], output_filename: str = "pitch_deck.pptx"):
    prs = Presentation()
    
    def add_slide(title_text, content_text):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        # Title Formatting
        title = slide.shapes.title
        title.text = title_text
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.size = Pt(28)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(0, 51, 102)
            paragraph.alignment = PP_ALIGN.LEFT
            
        # Body Formatting
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        p = tf.add_paragraph()
        p.text = content_text
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(50, 50, 50)

    # 1. Title Slide
    startup_name = state.get("startup_context", {}).get("startup_name", "Startup Pitch Deck")
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = startup_name
    subtitle.text = "AI-Generated Startup Accelerator Report"

    # 2. Problem & Solution
    startup_context = state.get("startup_context", {})
    add_slide("Problem & Solution", f"Problem:\n{startup_context.get('problem_statement', 'N/A')}\n\nSolution:\n{startup_context.get('solution_description', 'N/A')}")

    # 3. Market Intelligence
    market_int = state.get("market_intelligence", {}).get("market_analysis", "N/A")
    add_slide("Market Intelligence", market_int[:500] + "..." if len(market_int) > 500 else market_int)

    # 4. Business Validation
    biz_val = state.get("business_validation", {}).get("evaluation", "N/A")
    add_slide("Business Validation", biz_val[:500] + "..." if len(biz_val) > 500 else biz_val)

    # 5. Executive Strategy
    strategy = state.get("strategy_output", {}).get("strategy", "N/A")
    add_slide("Executive Strategy", strategy[:500] + "..." if len(strategy) > 500 else strategy)

    # Ensure reports directory exists
    reports_dir = os.path.join(os.getcwd(), "reports")
    os.makedirs(reports_dir, exist_ok=True)
    
    file_path = os.path.join(reports_dir, output_filename)
    prs.save(file_path)
    return file_path
